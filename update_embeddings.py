from drive_auth_test11 import authenticate_drive, fetch_documents
from sentence_transformers import SentenceTransformer
import numpy as np
import pickle
import os
from docx import Document
from pptx import Presentation

def extract_text_docx(file_path):
    """Extract all text from a .docx file."""
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_pptx(file_path):
    """Extract all text from a .pptx file."""
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

def chunk_text(text, max_tokens=200):
    """Split text into chunks of max_tokens words."""
    sentences = text.split('. ')  # Basic sentence split; consider nltk for better splits
    chunks = []
    chunk = ""
    for sentence in sentences:
        # +1 because sentence will add one period
        if len((chunk + sentence).split()) > max_tokens:
            chunks.append(chunk.strip())
            chunk = sentence + ". "
        else:
            chunk += sentence + ". "
    if chunk:
        chunks.append(chunk.strip())
    return chunks

print("🔄 Fetching documents from Google Drive...")
docs, sources_raw, file_ids_raw, file_paths = fetch_documents(authenticate_drive)

print("✂️ Extracting text, chunking and labeling...")

text_chunks = []
sources = []
file_ids = []

for doc, source, file_id, file_path in zip(docs, sources_raw, file_ids_raw, file_paths):
    ext = os.path.splitext(file_path or "")[1].lower()  # Handle None paths safely

    if ext == ".docx" and file_path and os.path.exists(file_path):
        text = extract_text_docx(file_path)
    elif ext == ".pptx" and file_path and os.path.exists(file_path):
        text = extract_text_pptx(file_path)
    else:
        # fallback to raw text extracted from Google Docs or PDFs
        text = doc

    chunks = chunk_text(text)

    # Add file name and folder name as searchable chunks for better matching
    if " / " in source:
        folder_name, file_name = source.split(" / ", 1)
    else:
        folder_name, file_name = "", source

    # Append filenames and folder names as chunks
    chunks.append(file_name)
    if folder_name:
        chunks.append(folder_name)

    text_chunks.extend(chunks)
    sources.extend([source] * len(chunks))
    file_ids.extend([file_id] * len(chunks))

print("🔍 Embedding chunks...")
model = SentenceTransformer('all-MiniLM-L6-v2')
embeddings = model.encode(text_chunks, show_progress_bar=True).astype('float32')

print("💾 Saving data...")
os.makedirs('data', exist_ok=True)
np.save('data/embeddings.npy', embeddings)

with open('data/text_chunks.pkl', 'wb') as f:
    pickle.dump(text_chunks, f)

with open('data/sources.pkl', 'wb') as f:
    pickle.dump(sources, f)

with open('data/file_ids.pkl', 'wb') as f:
    pickle.dump(file_ids, f)

print("✅ Embedding update complete!")
